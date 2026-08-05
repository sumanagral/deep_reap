#!/usr/bin/env python3
"""Build a DeepREAP conference-style PowerPoint from paper figures."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "paper" / "figures"
OUT = ROOT / "paper" / "DeepREAP_Presentation.pptx"

# Academic navy / warm sand palette (avoid purple-default look)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x68)
ACCENT = RGBColor(0xC4, 0x5C, 0x26)  # terracotta-adjacent but used sparingly
TEAL = RGBColor(0x1F, 0x6F, 0x6B)
LIGHT = RGBColor(0xF7, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD0, 0xCB, 0xC0)


def _set_run(run, text, size=18, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                 color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    # Support multi-line text as separate paragraphs
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        run = p.add_run()
        _set_run(run, line, size=size, bold=bold, color=color, font=font)
    return box


def _blank_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)


def _bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _bar(slide, color=NAVY, height=Inches(0.12)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _footer(slide, page, total=12):
    _add_textbox(
        slide, Inches(0.5), Inches(7.15), Inches(10), Inches(0.3),
        "DeepREAP  ·  SJSU CMPE", size=11, color=MUTED, font="Calibri"
    )
    _add_textbox(
        slide, Inches(11.5), Inches(7.15), Inches(1.3), Inches(0.3),
        f"{page} / {total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT, font="Calibri"
    )


def _title_block(slide, title, subtitle=None):
    _bar(slide)
    _add_textbox(
        slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55),
        title, size=28, bold=True, color=NAVY, font="Georgia"
    )
    if subtitle:
        _add_textbox(
            slide, Inches(0.55), Inches(0.9), Inches(12.2), Inches(0.4),
            subtitle, size=14, color=MUTED, font="Calibri"
        )
    # accent rule under title
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.25), Inches(1.4), Inches(0.04)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def _bullets(slide, left, top, width, height, items, *, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        run = p.add_run()
        _set_run(run, "•  " + item, size=size, color=INK, font="Calibri")
    return box


def _stat_card(slide, left, top, width, height, value, label, accent=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RULE
    card.line.width = Pt(1)
    _add_textbox(
        slide, left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.55),
        value, size=26, bold=True, color=accent, align=PP_ALIGN.CENTER, font="Georgia"
    )
    _add_textbox(
        slide, left + Inches(0.1), top + Inches(0.7), width - Inches(0.2), Inches(0.55),
        label, size=12, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri"
    )


def _picture(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        _add_textbox(slide, left, top, Inches(4), Inches(0.4),
                     f"[missing: {path.name}]", size=12, color=ACCENT)
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    total = 12

    # ---- 1 Title ----
    s = _blank_slide(prs)
    _bg(s, NAVY)
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    _add_textbox(
        s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.4),
        "DeepREAP", size=48, bold=True, color=WHITE, font="Georgia"
    )
    _add_textbox(
        s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2),
        "Integrating Regressive Ensemble Demand Prediction\n"
        "with Deep Reinforcement Learning for Cloud Resource Allocation",
        size=22, color=RGBColor(0xD8, 0xDE, 0xE8), font="Calibri"
    )
    _add_textbox(
        s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.8),
        "Suma Nagral  ·  Raghav Sharma\nSan José State University",
        size=16, color=RGBColor(0xA8, 0xB4, 0xC8), font="Calibri"
    )
    _add_textbox(
        s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
        "Forecast-conditioned SJF imitation  ·  latency–throughput operating point",
        size=14, color=ACCENT, font="Calibri"
    )

    # ---- 2 Problem ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "The problem", "Cloud allocation under dynamic demand")
    _bullets(s, Inches(0.55), Inches(1.55), Inches(12), Inches(4.5), [
        "Cloud resource allocation trades per-job latency against cluster throughput.",
        "SMEs feel this without hyperscale slack or dedicated capacity-planning staff.",
        "Two ML families address it separately:",
        "    DRL schedulers (DeepRM / DeepRM_Plus) react to current cluster state.",
        "    Regression ensembles (REAP) forecast future demand from telemetry.",
        "They are complementary — reactive vs anticipatory — but seldom combined.",
    ], size=18)
    _footer(s, 2, total)

    # ---- 3 Gap & idea ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "The integration axis", "What the scheduler knows about the future")
    # two columns
    left = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.6), Inches(5.7), Inches(4.4)
    )
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = RULE
    right = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.05), Inches(1.6), Inches(5.7), Inches(4.4)
    )
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = RULE
    _add_textbox(s, Inches(0.8), Inches(1.85), Inches(5.2), Inches(0.4),
                 "Vanilla DRL", size=20, bold=True, color=NAVY, font="Georgia")
    _bullets(s, Inches(0.8), Inches(2.4), Inches(5.2), Inches(3.2), [
        "Observes current cluster load",
        "Sees the visible job queue",
        "No demand forecast in state",
        "Policy is purely reactive",
    ], size=16)
    _add_textbox(s, Inches(7.3), Inches(1.85), Inches(5.2), Inches(0.4),
                 "DeepREAP", size=20, bold=True, color=TEAL, font="Georgia")
    _bullets(s, Inches(7.3), Inches(2.4), Inches(5.2), Inches(3.2), [
        "Same cluster + job image",
        "Plus REAP forecast channels",
        "Fused into state — not a hand-coded priority score",
        "Policy learns how much to trust the forecast",
    ], size=16)
    _footer(s, 3, total)

    # ---- 4 RQs & contributions ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Research questions & contributions")
    _add_textbox(s, Inches(0.55), Inches(1.55), Inches(12), Inches(0.35),
                 "Research questions", size=16, bold=True, color=ACCENT, font="Calibri")
    _bullets(s, Inches(0.55), Inches(1.9), Inches(12), Inches(1.5), [
        "RQ1 — How can REAP forecasts enter DeepRM_Plus state without a hand-coded predictor–scheduler contract?",
        "RQ2 — What latency / completion / throughput operating point does that integration produce, and which training regime realises it?",
    ], size=16)
    _add_textbox(s, Inches(0.55), Inches(3.5), Inches(12), Inches(0.35),
                 "Contributions", size=16, bold=True, color=ACCENT, font="Calibri")
    _bullets(s, Inches(0.55), Inches(3.85), Inches(12), Inches(2.5), [
        "(i) Working DeepREAP system: REAP forecasts fused into a DeepRM-style CNN state, with REST + monitoring.",
        "(ii) Twin evaluation: forecast channels are the only difference → slowdown 8.27 / completion 19.04 vs 19.30 / 43.41.",
        "(iii) Practical recipe: SJF imitation with forecast-augmented state; short PPO does not improve it here.",
    ], size=16)
    _footer(s, 4, total)

    # ---- 5 Pipeline ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Pipeline", "Four stages; forecasts fused into state before the policy")
    stages = [
        ("1. Workload", "60-day telemetry\n3,000-job Poisson trace"),
        ("2. REAP", "GA features →\n5 regressors → inv-MSE"),
        ("3. State fusion", "Cluster + jobs\n⊕ CPU/mem forecasts"),
        ("4. CNN policy", "SJF imitation →\nPPO (15 updates)"),
    ]
    for i, (title, body) in enumerate(stages):
        x = Inches(0.55 + i * 3.15)
        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.7), Inches(2.9), Inches(2.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RULE
        top_bar = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(1.7), Inches(2.9), Inches(0.12)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = TEAL if i != 2 else ACCENT
        top_bar.line.fill.background()
        _add_textbox(s, x + Inches(0.15), Inches(2.0), Inches(2.6), Inches(0.5),
                     title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Georgia")
        _add_textbox(s, x + Inches(0.15), Inches(2.65), Inches(2.6), Inches(1.3),
                     body, size=14, color=MUTED, align=PP_ALIGN.CENTER, font="Calibri")
        if i < 3:
            _add_textbox(s, x + Inches(2.85), Inches(2.7), Inches(0.35), Inches(0.4),
                         "→", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _bullets(s, Inches(0.55), Inches(4.6), Inches(12), Inches(2), [
        "Design choice: fuse forecasts into the state image — no hand-coded priority contract.",
        "Twin protocol: DeepREAP vs DeepRM_Plus differ only in those forecast channels.",
        "Environment: 2 resources × capacity 10, horizon T=20, M=5 visible jobs.",
    ], size=16)
    _footer(s, 5, total)

    # ---- 6 REAP accuracy ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "REAP prediction accuracy", "60-day synthetic trace — ensemble vs base regressors")
    _picture(s, FIG / "reap_quality_cpu_load.png", Inches(0.4), Inches(1.5), width=Inches(6.2))
    _picture(s, FIG / "reap_quality_memory_usage.png", Inches(6.7), Inches(1.5), width=Inches(6.2))
    _add_textbox(
        s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.6),
        "CPU: ensemble MSE 8.07 — best of all models.  Memory: competitive (25.04) but SVR alone wins (24.49).",
        size=14, color=MUTED, font="Calibri"
    )
    _footer(s, 6, total)

    # ---- 7 Ensemble weights ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Why memory is competitive, not dominant", "Inverse-MSE weights stay nearly flat")
    _picture(s, FIG / "reap_weights_cpu_load.png", Inches(0.5), Inches(1.55), width=Inches(6.0))
    _picture(s, FIG / "reap_weights_memory_usage.png", Inches(6.8), Inches(1.55), width=Inches(6.0))
    _bullets(s, Inches(0.55), Inches(5.5), Inches(12), Inches(1.4), [
        "Weights span only ~0.13–0.23 → weaker models keep non-trivial mass.",
        "Reasonable default when the best base model is unknown; sharper weighting is a natural next step.",
    ], size=16)
    _footer(s, 7, total)

    # ---- 8 Key scheduling result ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Key result: twin ablation", "Forecast channels are the only manipulated factor")
    _stat_card(s, Inches(0.55), Inches(1.55), Inches(3.9), Inches(1.5),
               "8.27", "DeepREAP slowdown\n(imitation)", TEAL)
    _stat_card(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(1.5),
               "19.30 → 8.27", "Slowdown with forecasts\nvs DeepRM_Plus twin", ACCENT)
    _stat_card(s, Inches(8.85), Inches(1.55), Inches(3.9), Inches(1.5),
               "19.04", "Avg completion (slots)\nvs 43.41 twin", TEAL)
    _add_textbox(s, Inches(0.55), Inches(3.3), Inches(12), Inches(0.35),
                 "Shared job trace (seed 123), 201-step budget, shared training recipe",
                 size=13, color=MUTED, font="Calibri")
    # mini table as bullets
    rows = [
        ("FIFO / SJF / Packer", "slowdown 18.76–20.20", "completion ~48–50"),
        ("DeepRM_Plus (imitation)", "19.30", "43.41"),
        ("DeepREAP (imitation)", "8.27  ← best", "19.04  ← best"),
        ("DeepREAP (PPO, 15 upd.)", "18.12", "42.05"),
    ]
    y = Inches(3.75)
    for name, a, b in rows:
        _add_textbox(s, Inches(0.7), y, Inches(4.5), Inches(0.35), name, size=15, bold=True, color=NAVY)
        _add_textbox(s, Inches(5.4), y, Inches(3.5), Inches(0.35), a, size=15, color=INK)
        _add_textbox(s, Inches(9.0), y, Inches(3.5), Inches(0.35), b, size=15, color=INK)
        y += Inches(0.42)
    _footer(s, 8, total)

    # ---- 9 Scheduling figures ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Scheduling metrics", "Lower is better — DeepREAP-imitation leads on both latency metrics")
    _picture(s, FIG / "avg_slowdown.png", Inches(0.4), Inches(1.5), width=Inches(6.2))
    _picture(s, FIG / "avg_completion.png", Inches(6.7), Inches(1.5), width=Inches(6.2))
    _add_textbox(
        s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.55),
        "Heuristics cluster ~18–20 slowdown; forecast-conditioned imitation is the clear latency winner.",
        size=14, color=MUTED, font="Calibri"
    )
    _footer(s, 9, total)

    # ---- 10 Latency–throughput ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Latency–throughput operating point", "Not a claim of higher work-conserving throughput")
    _stat_card(s, Inches(0.55), Inches(1.6), Inches(4.0), Inches(1.6),
               "25 jobs", "DeepREAP completed\nin 201 steps", ACCENT)
    _stat_card(s, Inches(4.8), Inches(1.6), Inches(4.0), Inches(1.6),
               "52–54", "Heuristic n_done\n(same budget)", TEAL)
    _stat_card(s, Inches(9.05), Inches(1.6), Inches(3.7), Inches(1.6),
               "−446.4", "Worst total reward\n(backlog still charged)", MUTED)
    _bullets(s, Inches(0.55), Inches(3.6), Inches(12), Inches(3), [
        "Policy defers long jobs into windows where REAP anticipates spare headroom.",
        "Finished jobs are much faster (slowdown / completion); raw throughput drops.",
        "Reward sums over all resident jobs — including backlog — so it is not a proxy for scheduling quality alone.",
        "Dollar cost = a rate card on this table; we leave the rate unspecified.",
    ], size=16)
    _footer(s, 10, total)

    # ---- 11 Training dynamics ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Training dynamics", "Short PPO does not improve on imitation at this budget")
    _picture(s, FIG / "learning_curve_deeprm_plus.png", Inches(0.4), Inches(1.45), width=Inches(6.2))
    _picture(s, FIG / "learning_curve_deepreap.png", Inches(6.7), Inches(1.45), width=Inches(6.2))
    _bullets(s, Inches(0.55), Inches(5.55), Inches(12), Inches(1.4), [
        "15 PPO updates (~1.5×10⁴ transitions): reward oscillates; no clear upward trend.",
        "Reported recipe: forecast-conditioned SJF imitation. Whether larger PPO helps remains open.",
    ], size=15)
    _footer(s, 11, total)

    # ---- 12 Takeaways ----
    s = _blank_slide(prs)
    _bg(s)
    _title_block(s, "Takeaways", "Narrow, defensible claim")
    _bullets(s, Inches(0.55), Inches(1.55), Inches(12), Inches(3.8), [
        "DeepREAP = state-fusion of REAP forecasts into a DeepRM-style scheduler.",
        "Under controlled synthetic short-horizon eval, SJF imitation with forecasts reaches slowdown 8.27 and completion 19.04 — vs 19.30 / 43.41 for the no-forecast twin.",
        "That win is a latency–throughput operating point (fewer jobs finished under a fixed step budget).",
        "REAP ensemble wins on CPU load; competitive but not dominant on memory.",
        "Limitations: synthetic trace, fixed seed/trace, fixed horizon — public traces and longer drains are next.",
    ], size=17)
    box = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    _add_textbox(
        s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.7),
        "Not “predictive DRL always wins” — a quantified forecast-in-state operating point under a fixed protocol.",
        size=16, bold=False, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri"
    )
    _footer(s, 12, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
